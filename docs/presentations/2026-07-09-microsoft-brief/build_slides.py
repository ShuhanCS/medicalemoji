"""Add three slides to the Emoji 2026 Brief deck, matching the existing phone-frame style.

Reads the source deck for everything it needs: the phone-frame chrome comes off the
"Request" slide, the 14 candidate emoji come off "Our New Candidates".
"""
import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SRC = r"C:\Users\Shuha\Downloads\Emoiji 2026 Brief.pptx"
OUT = "Emoji-2026-Brief-v2.pptx"
PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE

INK = RGBColor(0x48, 0x4E, 0x56)
SLATE = RGBColor(0x65, 0x6F, 0x7F)
MID = RGBColor(0xB7, 0xBD, 0xC6)
LIGHT = RGBColor(0xF3, 0xF3, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1E, 0x7A, 0x50)
RED = RGBColor(0xA3, 0x30, 0x2B)
AMBER = RGBColor(0x8A, 0x5E, 0x0F)
BLUE = RGBColor(0x2F, 0x5D, 0x8C)

HEAD = "Century Gothic"
BODY = "Calibri"

prs = Presentation(SRC)
BLANK = prs.slide_layouts[10]
slides = list(prs.slides)


def find_slide(title):
    """Locate a source slide by its title text, so adding slides upstream can't misalign us."""
    for s in slides:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() == title:
                return s
    raise LookupError(f"no slide titled {title!r} in {SRC}")


# Reusable chrome, lifted from the Request slide.
request_slide = find_slide("Request")
chrome = {}
for sh in request_slide.shapes:
    if sh.shape_type != PICTURE:
        continue
    w = Emu(sh.width).inches
    if w > 8:
        chrome["bg"] = sh.image.blob
    elif 0.9 < w < 1.0:
        chrome["logo"] = sh.image.blob
    elif w < 0.3:
        chrome["pin"] = sh.image.blob

# The 14 candidate emoji, with the native size they were placed at (for aspect ratio).
# Ordered the way they sit on the slide: top-to-bottom, then left-to-right.
_cands = [
    (Emu(sh.top).inches, Emu(sh.left).inches, Emu(sh.width).inches, Emu(sh.height).inches, sh)
    for sh in find_slide("Our New Candidates").shapes
    if sh.shape_type == PICTURE and Emu(sh.width).inches < 8 and Emu(sh.top).inches >= 1.0
]
_cands.sort(key=lambda r: (round(r[0] * 2) / 2, r[1]))
CANDIDATES = [(sh.image.blob, w, h) for _, _, w, h, sh in _cands]
assert len(CANDIDATES) == 14, f"expected 14 candidate emoji, found {len(CANDIDATES)}"

boxes = []  # geometry log for QA


def log(sl, name, l, t, w, h):
    boxes.append((sl, name, l, t, l + w, t + h))


def txbox(slide, sl, name, l, t, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    log(sl, name, l, t, w, h)
    return tf


def run(para, text, *, font=BODY, size=10, bold=False, color=INK, italic=False):
    r = para.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def lines(tf, text, *, align=None, **kw):
    """Render a \\n-separated string as one paragraph per line."""
    for j, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        run(para, line, **kw)


def card(slide, sl, name, l, t, w, h, *, fill=WHITE, line=MID, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    log(sl, name, l, t, w, h)
    return sh


def topbar(slide, l, t, w, color, inset=0.10):
    """Accent bar inset so it never pokes out of a rounded corner."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l + inset), Inches(t + 0.014), Inches(w - 2 * inset), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def base(title, sl):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(io.BytesIO(chrome["bg"]), 0, 0, Inches(10), Inches(5.625))
    if "pin" in chrome:
        s.shapes.add_picture(io.BytesIO(chrome["pin"]), Inches(0.27), Inches(0.59), Inches(0.17), Inches(0.29))
    if "logo" in chrome:
        s.shapes.add_picture(io.BytesIO(chrome["logo"]), Inches(8.74), Inches(0.26), Inches(0.94), Inches(0.17))
    tf = txbox(s, sl, "title", 0.60, 0.38, 8.00, 0.62)
    run(tf.paragraphs[0], title, font=HEAD, size=30, bold=True, color=INK)
    ff = txbox(s, sl, "footer", 0.18, 5.19, 1.60, 0.29)
    run(ff.paragraphs[0], "Shuhan He, MD", font=BODY, size=9, color=SLATE)
    return s


def sources(slide, sl, text):
    tf = txbox(slide, sl, "src", 1.95, 5.21, 7.35, 0.26, align=PP_ALIGN.RIGHT)
    run(tf.paragraphs[0], text, size=6.5, color=SLATE)


# ---------------------------------------------------------------- slide A
S = "A"
s = base("The 2021 JAMA set", S)
tf = txbox(s, S, "sub", 0.62, 1.08, 8.66, 0.34)
p = tf.paragraphs[0]
run(p, "Fourteen clinical concepts proposed for Unicode.  ", size=10.5, color=SLATE)
run(p, "Every one was declined.", size=10.5, bold=True, color=RED)

ORDER = [
    (5, "Kidneys"), (4, "Liver"), (2, "Stomach"), (3, "Spine"), (0, "Intestines"),
    (13, "White blood cell"), (12, "ECG"), (7, "Blood bag"), (8, "IV bag"), (6, "Pill pack"),
    (11, "Pill box"), (1, "Leg cast"), (9, "CT scan"), (10, "Weight scale"),
]
COLS, X0, CW = 5, 0.62, 1.744
ROWS_Y = [1.50, 2.58, 3.66]
IMG_H, CAP_H, CAP_GAP = 0.78, 0.22, 0.04
# Match apparent size by area, not height: a tall thin spine and a wide flat ECG
# both look wrong when normalised on a single dimension.
TARGET_AREA, MAX_W = 0.55, 1.12

for i, (idx, label) in enumerate(ORDER):
    c, r = i % COLS, i // COLS
    n_in_row = min(COLS, len(ORDER) - r * COLS)
    row_shift = (COLS - n_in_row) * CW / 2   # centre a short final row
    cx = X0 + c * CW + row_shift
    cy = ROWS_Y[r]
    blob, w0, h0 = CANDIDATES[idx]
    sc = min((TARGET_AREA / (w0 * h0)) ** 0.5, IMG_H / h0, MAX_W / w0)
    iw, ih = w0 * sc, h0 * sc
    s.shapes.add_picture(
        io.BytesIO(blob),
        Inches(cx + (CW - iw) / 2), Inches(cy + (IMG_H - ih)),
        Inches(iw), Inches(ih),
    )
    ctf = txbox(s, S, f"cap{i}", cx, cy + IMG_H + CAP_GAP, CW, CAP_H, align=PP_ALIGN.CENTER)
    run(ctf.paragraphs[0], label, size=9, color=INK)

note = card(s, S, "note", 0.62, 4.76, 8.66, 0.36, fill=LIGHT, line=LIGHT)
ntf = note.text_frame
ntf.word_wrap = True
ntf.margin_left = ntf.margin_right = Inches(0.14)
ntf.margin_top = ntf.margin_bottom = 0
ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
np = ntf.paragraphs[0]
np.alignment = PP_ALIGN.LEFT  # autoshapes inherit centred text
run(np, "Two earlier proposals did ship: ", size=9.5, color=SLATE)
run(np, "the anatomical heart and lungs, Emoji 13.0 (2020).", size=9.5, bold=True, color=GREEN)

sources(s, S, "Lai D, Lee J, He S. Emoji for the Medical Community — Challenges and Opportunities. JAMA. 2021;326(9):795-796.  ·  Status: unicode.org/emoji/emoji-proposals-status.html")

# ---------------------------------------------------------------- slide B
S = "B"
s = base("Every resubmission reset the clock", S)
tf = txbox(s, S, "sub", 0.62, 1.08, 8.66, 0.32)
p = tf.paragraphs[0]
run(p, "Unicode bars a declined emoji from re-review for ", size=10.5, color=SLATE)
run(p, "four years", size=10.5, bold=True, color=INK)
run(p, ", counted from the decline.", size=10.5, color=SLATE)

# timeline rail
RAIL_Y, LX, SPAN = 2.06, 0.86, 8.20
rail = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LX), Inches(RAIL_Y), Inches(SPAN), Inches(0.020))
rail.fill.solid()
rail.fill.fore_color.rgb = MID
rail.line.fill.background()
rail.shadow.inherit = False
log(S, "rail", LX, RAIL_Y, SPAN, 0.02)

EVENTS = [
    ("2019", "Heart + lungs\nsubmitted", SLATE),
    ("2020", "Both ship in\nEmoji 13.0", GREEN),
    ("2021", "First wave declined;\nJAMA makes the case", RED),
    ("2022", "Kidney, stomach,\nliver declined", AMBER),   # ties to the knife-edge card
    ("2024", "Spine, intestines,\nECG declined", RED),
    ("2026", "Window shuts\nJuly 31", BLUE),
]
N = len(EVENTS)
for i, (yr, cap, col) in enumerate(EVENTS):
    cx = LX + (SPAN / (N - 1)) * i
    d = 0.13 if col is not BLUE else 0.17
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(RAIL_Y + 0.01 - d / 2), Inches(d), Inches(d))
    dot.fill.solid()
    dot.fill.fore_color.rgb = col
    dot.line.color.rgb = WHITE
    dot.line.width = Pt(1.25)
    dot.shadow.inherit = False
    ytf = txbox(s, S, f"yr{i}", cx - 0.70, RAIL_Y - 0.46, 1.40, 0.28, align=PP_ALIGN.CENTER)
    run(ytf.paragraphs[0], yr, font=HEAD, size=13, bold=True, color=col)
    ctf = txbox(s, S, f"ev{i}", cx - 0.70, RAIL_Y + 0.20, 1.40, 0.52, align=PP_ALIGN.CENTER)
    lines(ctf, cap, align=PP_ALIGN.CENTER, size=8, color=SLATE)

# eligibility columns
CY, CH = 3.06, 1.58
COLS_B = [
    ("Clear to file now", GREEN, "8 concepts",
     "White blood cell · Pill pack · Pill box\nBlood bag · Leg cast · IV bag\nCT scan · Weight scale",
     "Last declined in the 2021 cycle. The four-year bar has lapsed."),
    ("On a knife edge", AMBER, "3 concepts",
     "Kidney · Stomach · Liver",
     "Submitted July 2022, declined that November. From the decline, the bar lifts Nov 2026 — after the window shuts. From submission, they clear by 12, 3 and 1 day."),
    ("Barred", RED, "3 concepts",
     "Spine · Intestines · ECG",
     "Submitted days after the 2024 window opened, declined that cycle. The bar lifts around Nov 2028."),
]
cw, gap = 2.78, 0.16
for i, (head, col, count, items, note_t) in enumerate(COLS_B):
    cx = 0.62 + i * (cw + gap)
    c = card(s, S, f"col{i}", cx, CY, cw, CH)
    topbar(s, cx, CY, cw, col)

    htf = txbox(s, S, f"h{i}", cx + 0.14, CY + 0.16, cw - 0.28, 0.24)
    hp = htf.paragraphs[0]
    run(hp, head, font=HEAD, size=11.5, bold=True, color=col)
    run(hp, f"   {count}", size=8.5, color=MID)

    itf = txbox(s, S, f"i{i}", cx + 0.14, CY + 0.44, cw - 0.28, 0.52)
    lines(itf, items, size=8.8, bold=True, color=INK)

    ntf = txbox(s, S, f"n{i}", cx + 0.14, CY + 0.96, cw - 0.28, 0.52)
    run(ntf.paragraphs[0], note_t, size=7.4, color=SLATE)

sources(s, S, "Unicode publishes submission dates, not decline dates; declines issue at cycle end, submitters notified by Nov 30.  Sources: unicode.org/emoji/emoji-proposals-status.html · unicode.org/emoji/proposals.html")

# ---------------------------------------------------------------- slide C
S = "C"
s = base("Where emoji decisions get made", S)
tf = txbox(s, S, "sub", 0.62, 1.08, 8.66, 0.32)
p = tf.paragraphs[0]
run(p, "The subcommittee recommends. The Technical Committee votes. ", size=10.5, color=SLATE)
run(p, "Microsoft votes in that room.", size=10.5, bold=True, color=BLUE)

CHAIN = [
    ("Proposal", "Anyone may file\nduring the window", False),
    ("Emoji Standard &\nResearch WG", "formerly the Emoji\nSubcommittee (ESC)\nrecommends only", False),
    ("Unicode Technical\nCommittee", "holds the vote\nFull = 1 · Supporting = ½", True),
    ("Release + fonts", "Unicode version, then\nSegoe UI Emoji ships it", False),
]
BW, BG, BY, BH = 1.95, 0.28, 1.52, 1.24
for i, (head, sub, ms) in enumerate(CHAIN):
    bx = 0.62 + i * (BW + BG)
    # No accent bars here: one blue border marks the only room that votes.
    c = card(s, S, f"box{i}", bx, BY, BW, BH, fill=WHITE, line=BLUE if ms else MID)
    if ms:
        c.line.width = Pt(1.75)

    htf = txbox(s, S, f"ch{i}", bx + 0.12, BY + 0.20, BW - 0.24, 0.44, align=PP_ALIGN.CENTER)
    lines(htf, head, align=PP_ALIGN.CENTER, font=HEAD, size=10, bold=True, color=BLUE if ms else INK)

    stf = txbox(s, S, f"cs{i}", bx + 0.10, BY + 0.66, BW - 0.20, 0.50, align=PP_ALIGN.CENTER)
    lines(stf, sub, align=PP_ALIGN.CENTER, size=7.6, color=SLATE)

    if i < len(CHAIN) - 1:
        ax = bx + BW + (BG - 0.16) / 2
        ar = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(ax), Inches(BY + BH / 2 - 0.10), Inches(0.16), Inches(0.20))
        ar.rotation = 90
        ar.fill.solid()
        ar.fill.fore_color.rgb = MID
        ar.line.fill.background()
        ar.shadow.inherit = False

# Microsoft card
MY, MH = 3.06, 1.62
mc = card(s, S, "ms", 0.62, MY, 5.34, MH, fill=WHITE, line=BLUE)
mc.line.width = Pt(1.75)
htf = txbox(s, S, "mh", 0.78, MY + 0.16, 5.02, 0.26)
run(htf.paragraphs[0], "Microsoft's standing in the Consortium", font=HEAD, size=12, bold=True, color=BLUE)

BULLETS = [
    ("Full member.", " One of the few tiers that carries a UTC vote."),
    ("Board seat.", " Vishal Chowdhary, VP of Science, Microsoft 365 Copilot."),
    ("Vendor.", " Segoe UI Emoji puts the glyph on every Windows device."),
]
for i, (b, rest) in enumerate(BULLETS):
    y = MY + 0.52 + i * 0.34
    btf = txbox(s, S, f"mb{i}", 0.92, y, 4.88, 0.30)
    bp = btf.paragraphs[0]
    run(bp, "—  ", size=9.5, color=BLUE)
    run(bp, b, size=9.5, bold=True, color=INK)
    run(bp, rest, size=9.5, color=SLATE)

# Calendar card
cc = card(s, S, "cal", 6.16, MY, 3.12, MH, fill=LIGHT, line=MID)
ctf = txbox(s, S, "calh", 6.32, MY + 0.16, 2.80, 0.24)
run(ctf.paragraphs[0], "The 2026 cycle", font=HEAD, size=12, bold=True, color=INK)
CAL = [("Opened", "April 2, 2026"), ("Closes", "July 31, 2026"), ("Decisions", "by Nov 30, 2026")]
for i, (k, v) in enumerate(CAL):
    y = MY + 0.52 + i * 0.30
    ktf = txbox(s, S, f"ck{i}", 6.32, y, 0.94, 0.26)
    run(ktf.paragraphs[0], k, size=9, color=SLATE)
    vtf = txbox(s, S, f"cv{i}", 7.26, y, 1.86, 0.26)
    run(vtf.paragraphs[0], v, size=9, bold=True, color=INK)
dtf = txbox(s, S, "cd", 6.32, MY + 1.42, 2.80, 0.22)
run(dtf.paragraphs[0], "Declined emoji: four-year bar", size=7.8, italic=True, color=RED)

sources(s, S, "unicode.org/consortium/utc.html  ·  unicode.org/consortium/directors.html  ·  unicode.org/emoji/proposals.html")

# ---------------------------------------------------------------- reorder: new slides before "Request"
# Push the Request slide to the end so the governance infographic lands just before the ask.
sldIdLst = prs.slides._sldIdLst
node = list(sldIdLst)[slides.index(request_slide)]
sldIdLst.remove(node)
sldIdLst.append(node)

prs.save(OUT)
print("saved", OUT, "| slides:", len(sldIdLst))

# ---------------------------------------------------------------- geometry QA
print("\n--- geometry QA ---")
SW, SH = 10.0, 5.625
bad = 0
for sl, name, l, t, r, b in boxes:
    if l < 0.14 or t < 0.14 or r > SW - 0.14 or b > SH - 0.10:
        print(f"  OUT-OF-BOUNDS {sl}/{name}: ({l:.2f},{t:.2f})-({r:.2f},{b:.2f})")
        bad += 1
by_slide = {}
for box in boxes:
    by_slide.setdefault(box[0], []).append(box)
SKIP = {"title", "src", "footer"}
for sl, items in by_slide.items():
    for i in range(len(items)):
        for j in range(i + 1, len(items)):
            _, n1, l1, t1, r1, b1 = items[i]
            _, n2, l2, t2, r2, b2 = items[j]
            if n1 in SKIP or n2 in SKIP:
                continue
            # cards legitimately contain their own text
            if n1.startswith(("col", "box", "ms", "cal", "note")) or n2.startswith(("col", "box", "ms", "cal", "note")):
                continue
            ox = min(r1, r2) - max(l1, l2)
            oy = min(b1, b2) - max(t1, t2)
            if ox > 0.01 and oy > 0.01:
                print(f"  OVERLAP {sl}: {n1} x {n2}  ({ox:.2f} x {oy:.2f} in)")
                bad += 1
print(f"--- {bad} geometry issues, {len(boxes)} boxes checked ---")
