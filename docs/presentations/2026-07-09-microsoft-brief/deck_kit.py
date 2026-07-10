"""Shared drawing helpers for the Emoji 2026 Brief deck.

Everything here matches the source deck's phone-frame style: Century Gothic titles in
`INK`, a `MID` frame, and chrome (background, status-bar icons, title pin) lifted off an
existing slide rather than re-drawn.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import io

PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE
SLIDE_W, SLIDE_H = 10.0, 5.625

INK = RGBColor(0x48, 0x4E, 0x56)
SLATE = RGBColor(0x65, 0x6F, 0x7F)
MID = RGBColor(0xB7, 0xBD, 0xC6)
LIGHT = RGBColor(0xF3, 0xF3, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1E, 0x7A, 0x50)
RED = RGBColor(0xA3, 0x30, 0x2B)
AMBER = RGBColor(0x8A, 0x5E, 0x0F)
BLUE = RGBColor(0x2F, 0x5D, 0x8C)

HEAD = "Century Gothic"
BODY = "Calibri"


def run(para, text, *, font=BODY, size=10, bold=False, color=INK, italic=False):
    r = para.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def lines(tf, text, *, align=None, **kw):
    """Render a \\n-separated string as one paragraph per line."""
    for j, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        run(para, line, **kw)


class Deck:
    """A source deck plus the geometry log used by the QA pass."""

    def __init__(self, src, chrome_from="Request"):
        self.src = src
        self.prs = Presentation(src)
        self.blank = self.prs.slide_layouts[10]
        self.boxes = []
        self.chrome = self._lift_chrome(self.find_slide(chrome_from))

    # -- source deck ------------------------------------------------------
    @property
    def slides(self):
        return list(self.prs.slides)

    def find_slide(self, title):
        """Locate a slide by title text, so adding slides upstream can't misalign us."""
        for s in self.slides:
            for sh in s.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip() == title:
                    return s
        raise LookupError(f"no slide titled {title!r} in {self.src}")

    @staticmethod
    def _lift_chrome(slide):
        chrome = {}
        for sh in slide.shapes:
            if sh.shape_type != PICTURE:
                continue
            w = Emu(sh.width).inches
            if w > 8:
                chrome["bg"] = sh.image.blob
            elif 0.9 < w < 1.0:
                chrome["logo"] = sh.image.blob
            elif w < 0.3:
                chrome["pin"] = sh.image.blob
        return chrome

    def drop_slide(self, slide):
        """Remove a slide from the deck, dropping its relationship."""
        idx = self.slides.index(slide)
        sld_id_lst = self.prs.slides._sldIdLst
        node = list(sld_id_lst)[idx]
        self.prs.part.drop_rel(node.rId)
        sld_id_lst.remove(node)

    def move_slide(self, slide, to_index):
        sld_id_lst = self.prs.slides._sldIdLst
        node = list(sld_id_lst)[self.slides.index(slide)]
        sld_id_lst.remove(node)
        sld_id_lst.insert(to_index, node)

    def save(self, out):
        self.prs.save(out)

    # -- drawing ----------------------------------------------------------
    def log(self, sl, name, l, t, w, h):
        self.boxes.append((sl, name, l, t, l + w, t + h))

    def base(self, title, sl):
        """A new blank slide carrying the deck's chrome, title and footer."""
        s = self.prs.slides.add_slide(self.blank)
        s.shapes.add_picture(io.BytesIO(self.chrome["bg"]), 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))
        if "pin" in self.chrome:
            s.shapes.add_picture(io.BytesIO(self.chrome["pin"]), Inches(0.27), Inches(0.59), Inches(0.17), Inches(0.29))
        if "logo" in self.chrome:
            s.shapes.add_picture(io.BytesIO(self.chrome["logo"]), Inches(8.74), Inches(0.26), Inches(0.94), Inches(0.17))
        tf = self.txbox(s, sl, "title", 0.60, 0.38, 8.00, 0.62)
        run(tf.paragraphs[0], title, font=HEAD, size=30, bold=True, color=INK)
        ff = self.txbox(s, sl, "footer", 0.18, 5.19, 1.60, 0.29)
        run(ff.paragraphs[0], "Shuhan He, MD", font=BODY, size=9, color=SLATE)
        return s

    def txbox(self, slide, sl, name, l, t, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        tf.paragraphs[0].alignment = align
        self.log(sl, name, l, t, w, h)
        return tf

    def card(self, slide, sl, name, l, t, w, h, *, fill=WHITE, line=MID, radius=0.06, width_pt=0.75):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.adjustments[0] = radius
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.color.rgb = line
        sh.line.width = Pt(width_pt)
        sh.shadow.inherit = False
        self.log(sl, name, l, t, w, h)
        return sh

    def topbar(self, slide, l, t, w, color, inset=0.10):
        """Accent bar inset so it never pokes out of a rounded corner."""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l + inset), Inches(t + 0.014), Inches(w - 2 * inset), Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.shadow.inherit = False

    def rule(self, slide, l, t, w, h, color=MID):
        """A hairline connector. Width or height of ~0.02in reads as a 1px line."""
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        sh.shadow.inherit = False

    def sources(self, slide, sl, text):
        tf = self.txbox(slide, sl, "src", 1.95, 5.21, 7.35, 0.26, align=PP_ALIGN.RIGHT)
        run(tf.paragraphs[0], text, size=6.5, color=SLATE)

    # -- QA ---------------------------------------------------------------
    def qa(self, skip=("title", "src", "footer"), containers=()):
        bad = 0
        for sl, name, l, t, r, b in self.boxes:
            if l < 0.14 or t < 0.14 or r > SLIDE_W - 0.14 or b > SLIDE_H - 0.10:
                print(f"  OUT-OF-BOUNDS {sl}/{name}: ({l:.2f},{t:.2f})-({r:.2f},{b:.2f})")
                bad += 1
        by_slide = {}
        for box in self.boxes:
            by_slide.setdefault(box[0], []).append(box)
        for sl, items in by_slide.items():
            for i in range(len(items)):
                for j in range(i + 1, len(items)):
                    _, n1, l1, t1, r1, b1 = items[i]
                    _, n2, l2, t2, r2, b2 = items[j]
                    if n1 in skip or n2 in skip:
                        continue
                    if n1.startswith(containers) or n2.startswith(containers):
                        continue  # cards legitimately contain their own text
                    ox, oy = min(r1, r2) - max(l1, l2), min(b1, b2) - max(t1, t2)
                    if ox > 0.01 and oy > 0.01:
                        print(f"  OVERLAP {sl}: {n1} x {n2}  ({ox:.2f} x {oy:.2f} in)")
                        bad += 1
        print(f"--- {bad} geometry issues, {len(self.boxes)} boxes checked ---")
        return bad
